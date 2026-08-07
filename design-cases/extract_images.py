#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""从 PPTX 每页提取最大图片 → slides/S###.jpg（1000px JPEG）"""
import os, subprocess, sys
from pptx import Presentation

BASE = '/Users/xiangzhong/Projects/personal-site/design-cases'
RAW = os.path.join(BASE, 'raw')
SLIDES = os.path.join(BASE, 'slides')
os.makedirs(RAW, exist_ok=True)
os.makedirs(SLIDES, exist_ok=True)

p = Presentation(os.path.join(BASE, 'zhouyu.pptx'))
ok, skip, err = [], [], []

for i, s in enumerate(p.slides, 1):
    best, best_area = None, 0
    for sh in s.shapes:
        if sh.shape_type == 13:  # PICTURE
            try:
                img = sh.image
                if img.content_type not in ('image/jpeg', 'image/png'):
                    continue
                area = (img.size[0] * img.size[1])
                if area > best_area:
                    best, best_area = img, area
            except Exception:
                continue
    if best is None:
        skip.append(i)
        continue
    ext = 'png' if best.content_type == 'image/png' else 'jpg'
    raw_path = os.path.join(RAW, f'S{i:03d}.{ext}')
    with open(raw_path, 'wb') as fh:
        fh.write(best.blob)
    # 压缩: 若 PNG 先转 JPEG, 统一 1000px q78
    out = os.path.join(SLIDES, f'S{i:03d}.jpg')
    if ext == 'png':
        subprocess.run(['sips', '-s', 'format', 'jpeg', '-s', 'formatOptions', '78',
                        '-Z', '1000', raw_path, '--out', out],
                       capture_output=True)
    else:
        subprocess.run(['sips', '-Z', '1000', '-s', 'format', 'jpeg',
                        '-s', 'formatOptions', '78', raw_path, '--out', out],
                       capture_output=True)
    if os.path.exists(out):
        ok.append(i)
    else:
        err.append(i)

print(f'成功: {len(ok)} 页, 无图跳过: {skip}, 失败: {err}')
total = sum(os.path.getsize(os.path.join(SLIDES, f'S{i:03d}.jpg')) for i in ok)
print(f'总大小: {total/1024/1024:.1f} MB')
